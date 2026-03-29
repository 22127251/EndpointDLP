"""
mitmproxy addon: file upload interceptor for Windows DLP.

Run with:
    mitmdump -s addon.py --listen-port 8080

Browser Channel Flow:
1. Intercept HTTP file upload (POST/PUT)
2. Extract text from text-based files (.txt, .csv, .md, etc.)
3. Chunk text into 500-word segments with 50-word overlap
4. Send chunks to QueueManager via named pipe
5. QueueManager analyzes and returns decision
6. Block upload if any chunk is BLOCKED
"""

import base64
import datetime
import json
import logging
import os
import tempfile
import threading
import time
from email.parser import BytesHeaderParser
from typing import List, Optional
from urllib.parse import urlparse, parse_qs

from mitmproxy import http

import pipe_client
from pipe_client import ChunkPayload, chunk_text
from config import Config, load_config

log = logging.getLogger(__name__)

_upload_lock = threading.Lock()
_cfg: Config = Config()

# Resumable upload tracking
_pending_resumable: dict = {}   # flow id → filename
_resumable_filenames: dict = {} # upload_id → filename
_blocked_url_cache: dict = {}   # url → expiry_time
_blocked_url_cache_lock = threading.Lock()
_BLOCK_CACHE_TTL = 60.0

# Track message decisions for browser uploads
_browser_decisions: dict = {}  # message_id -> "ALLOW" or "BLOCK"
_browser_decisions_lock = threading.Lock()


def load(loader):
    global _cfg
    # Use absolute path based on addon.py location
    import os
    addon_dir = os.path.dirname(os.path.abspath(__file__))
    config_path = os.path.join(addon_dir, "config.yaml")
    _cfg = load_config(config_path)
    tmp = _cfg.resolved_temp_dir()
    try:
        os.makedirs(tmp, exist_ok=True)
    except OSError as e:
        log.error("Cannot create temp dir %s: %s", tmp, e)
    log.info(
        "DLP addon loaded | pipe=%s timeout=%ss fail=%s chunk_size=%d overlap=%d",
        _cfg.pipe_name,
        _cfg.timeout_seconds,
        _cfg.fail_behavior,
        _cfg.chunk_size_words,
        _cfg.chunk_overlap_words,
    )


def requestheaders(flow: http.HTTPFlow) -> None:
    """Force full body buffering for potential uploads."""
    if flow.request.method not in ("POST", "PUT"):
        return

    host = flow.request.pretty_host.lower()
    if any(host == d or host.endswith("." + d) for d in _cfg.domain_blocklist):
        return

    content_type = flow.request.headers.get("content-type", "").lower()

    if "multipart/form-data" in content_type or "multipart/related" in content_type:
        log.debug("Force-buffering %s %s (content-type: %s)", flow.request.method, flow.request.pretty_url, content_type)
        flow.request.stream = False
        return

    # Force-buffer resumable uploads (Google Drive)
    if "uploadType=resumable" in flow.request.path.lower():
        log.debug("Force-buffering %s %s (resumable upload)", flow.request.method, flow.request.pretty_url)
        flow.request.stream = False
        return

    if any(kw in flow.request.path.lower() for kw in _cfg.upload_url_keywords):
        log.debug("Force-buffering %s %s (url keyword match)", flow.request.method, flow.request.pretty_url)
        flow.request.stream = False


def request(flow: http.HTTPFlow) -> None:
    # Handle batch resumable initiation
    if (flow.request.method == "PUT"
            and "multipart/mixed" in flow.request.headers.get("content-type", "").lower()
            and "batch" in flow.request.path.lower()):
        _track_resumable_initiation_batch(flow)
        return

    # Handle Google Drive resumable upload initiation (POST)
    if flow.request.method == "POST" and "uploadType=resumable" in flow.request.path:
        _track_resumable_upload_initiation(flow)
        return

    # Handle Google Drive resumable upload data (PUT)
    if flow.request.method == "PUT" and "uploadType=resumable" in flow.request.path:
        _handle_resumable_upload(flow)
        return

    # Debug: log all POST/PUT requests
    if flow.request.method in ("POST", "PUT"):
        log.debug("Checking %s %s | host=%s | content-type=%s",
                  flow.request.method, flow.request.pretty_url,
                  flow.request.pretty_host,
                  flow.request.headers.get("content-type", ""))

    if not _is_upload(flow):
        return

    log.info("Detected upload: %s %s", flow.request.method, flow.request.pretty_url)

    url = flow.request.pretty_url

    # Check if already blocked
    if _is_blocked_url(url):
        log.debug("BLOCK (cached) | %s", url[:80])
        flow.response = http.Response.make(403, b"Upload blocked by DLP policy.", {"Content-Type": "text/plain"})
        return

    body = flow.request.content
    content_type = flow.request.headers.get("content-type", "")

    if not body:
        log.warning("Empty body for upload %s %s", flow.request.method, flow.request.pretty_url)
        return

    # Extract filename and body
    if "multipart/related" in content_type.lower():
        filename, file_body, file_mime = _parse_multipart_related(body, content_type)
        if not filename:
            # Fallback: try to extract filename from URL
            filename = _extract_filename_from_url(flow.request.path)
            file_body = body
            file_mime = content_type.split(";")[0].strip().lower()
            if filename:
                log.info("Using fallback filename from URL: %s", filename)
        
        if not filename:
            filename = f"upload_{int(time.time())}"
            file_body = body
            file_mime = content_type.split(";")[0].strip().lower()
            log.info("Using generated filename: %s", filename)
    else:
        filename = _extract_filename(flow)
        file_body = body
        file_mime = content_type.split(";")[0].strip().lower()

    if not _matches_type_filter(filename, file_mime):
        log.debug("SKIP (type filter) | %s | %s", filename, file_mime)
        return

    # Save file to temp and send path to QueueManager for analysis
    _handle_upload_with_queue_analysis(flow, filename, file_body, file_mime, url)


def _handle_upload_with_queue_analysis(flow: http.HTTPFlow, filename: str, body: bytes, mime: str, url: str) -> None:
    """
    New workflow: Save file to temp, send path to QueueManager.
    QueueManager will: extract text → chunk → analyze → return decision.
    """
    try:
        temp_path = _write_temp_file(body, filename)
        log.info("Saved file to temp: %s (%d bytes)", temp_path, len(body))
    except OSError as e:
        log.error("Failed to write temp file for '%s': %s → fail_%s", filename, e, _cfg.fail_behavior)
        if not _cfg.fail_open():
            flow.response = http.Response.make(403, b"DLP Error: Cannot create temp file", {"Content-Type": "text/plain"})
        return

    # Send temp path to QueueManager for analysis
    payload = {
        "action": "analyze_file",
        "temp_path": temp_path,
        "filename": filename,
        "url": url,
        "method": flow.request.method,
        "content_type": flow.request.headers.get("content-type", ""),
        "mime_type": mime,
        "size_bytes": len(body),
        "timestamp": datetime.datetime.utcnow().isoformat() + "Z",
    }

    try:
        log.info("Sending file analysis request to pipe %s", _cfg.pipe_name)
        decision = pipe_client.send_and_receive(
            payload,
            _cfg.pipe_name,
            _cfg.timeout_seconds,
        )
        log.info("QueueManager decision: %s for %s", decision, filename)
    except Exception as e:
        log.warning("Pipe error: %s → fail_%s", e, _cfg.fail_behavior)
        decision = "ALLOW" if _cfg.fail_open() else "BLOCK"

    # Apply decision
    if decision == "BLOCK":
        _cache_blocked_url(url)
        log.info("BLOCK | %s | %d bytes | %s", filename, len(body), url)
        flow.response = http.Response.make(403, b"Upload blocked by DLP policy.", {"Content-Type": "text/plain"})
    else:
        log.info("ALLOW | %s | %d bytes | %s", filename, len(body), url)
        # File already uploaded (mitmproxy passed it through)
        # Decision is logged for audit


def _is_text_file(filename: str, mime: str) -> bool:
    """Check if file is suitable for text extraction and chunk analysis."""
    # Text files
    text_extensions = {'.txt', '.csv', '.md', '.json', '.xml', '.html', '.htm', '.log'}
    text_mimes = {'text/plain', 'text/csv', 'application/json', 'text/xml', 'text/html'}
    
    # Office documents
    office_extensions = {'.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx'}
    office_mimes = {
        'application/msword',
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        'application/vnd.ms-excel',
        'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        'application/vnd.ms-powerpoint',
        'application/vnd.openxmlformats-officedocument.presentationml.presentation',
    }
    
    # PDF
    pdf_extensions = {'.pdf'}
    pdf_mimes = {'application/pdf'}
    
    # Archives (will extract text files from inside)
    archive_extensions = {'.zip', '.tar', '.gz', '.rar', '.7z'}
    
    ext = os.path.splitext(filename)[1].lower()
    
    return (ext in text_extensions or mime in text_mimes or
            ext in office_extensions or mime in office_mimes or
            ext in pdf_extensions or mime in pdf_mimes or
            ext in archive_extensions)


def _extract_text_from_file(body: bytes, filename: str, mime: str) -> str:
    """Extract text content from uploaded file (supports text, Office, PDF, archives)."""
    ext = os.path.splitext(filename)[1].lower()
    
    # Try to decode as plain text first
    if ext in {'.txt', '.csv', '.md', '.json', '.xml', '.html', '.htm', '.log'}:
        return _extract_text_plain(body)
    
    # Office documents
    if ext in {'.doc', '.docx'}:
        return _extract_text_docx(body)
    if ext in {'.xls', '.xlsx'}:
        return _extract_text_xlsx(body)
    if ext in {'.ppt', '.pptx'}:
        return _extract_text_pptx(body)
    
    # PDF
    if ext == '.pdf':
        return _extract_text_pdf(body)
    
    # Archives
    if ext in {'.zip', '.tar', '.gz', '.rar', '.7z'}:
        return _extract_text_archive(body, filename)
    
    # Fallback: try plain text
    return _extract_text_plain(body)


def _extract_text_plain(body: bytes) -> str:
    """Extract text from plain text file."""
    encodings = ['utf-8', 'utf-8-sig', 'latin-1', 'cp1252']
    for encoding in encodings:
        try:
            return body.decode(encoding)
        except (UnicodeDecodeError, LookupError):
            continue
    return body.decode('utf-8', errors='replace')


def _extract_text_docx(body: bytes) -> str:
    """Extract text from DOCX file."""
    try:
        import docx
        from io import BytesIO
        doc = docx.Document(BytesIO(body))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return '\n'.join(paragraphs)
    except ImportError:
        log.warning("python-docx not installed, cannot extract text from DOCX")
        return ""
    except Exception as e:
        log.warning("Failed to extract text from DOCX: %s", e)
        return ""


def _extract_text_xlsx(body: bytes) -> str:
    """Extract text from XLSX file."""
    try:
        import openpyxl
        from io import BytesIO
        wb = openpyxl.load_workbook(BytesIO(body), read_only=True, data_only=True)
        texts = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_text = ' '.join(str(cell) if cell is not None else '' for cell in row if cell)
                if row_text.strip():
                    texts.append(row_text)
        return '\n'.join(texts)
    except ImportError:
        log.warning("openpyxl not installed, cannot extract text from XLSX")
        return ""
    except Exception as e:
        log.warning("Failed to extract text from XLSX: %s", e)
        return ""


def _extract_text_pptx(body: bytes) -> str:
    """Extract text from PPTX file."""
    try:
        import pptx
        from io import BytesIO
        prs = pptx.Presentation(BytesIO(body))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
        return '\n'.join(texts)
    except ImportError:
        log.warning("python-pptx not installed, cannot extract text from PPTX")
        return ""
    except Exception as e:
        log.warning("Failed to extract text from PPTX: %s", e)
        return ""


def _extract_text_pdf(body: bytes) -> str:
    """Extract text from PDF file."""
    try:
        import pypdf
        from io import BytesIO
        reader = pypdf.PdfReader(BytesIO(body))
        texts = []
        for page in reader.pages:
            text = page.extract_text()
            if text and text.strip():
                texts.append(text)
        return '\n'.join(texts)
    except ImportError:
        log.warning("pypdf not installed, cannot extract text from PDF")
        return ""
    except Exception as e:
        log.warning("Failed to extract text from PDF: %s", e)
        return ""


def _extract_text_archive(body: bytes, filename: str) -> str:
    """Extract text from archive files (zip, tar, etc.)."""
    import zipfile
    import tarfile
    import gzip
    
    ext = os.path.splitext(filename)[1].lower()
    texts = []
    
    try:
        if ext == '.zip':
            from io import BytesIO
            with zipfile.ZipFile(BytesIO(body), 'r') as zf:
                for name in zf.namelist():
                    if _is_text_file(name, ''):
                        try:
                            content = zf.read(name).decode('utf-8', errors='replace')
                            texts.append(f"=== {name} ===\n{content}")
                        except Exception:
                            pass
        
        elif ext == '.tar':
            from io import BytesIO
            with tarfile.open(fileobj=BytesIO(body), mode='r') as tf:
                for member in tf.getmembers():
                    if member.isfile() and _is_text_file(member.name, ''):
                        try:
                            content = tf.extractfile(member).read().decode('utf-8', errors='replace')
                            texts.append(f"=== {member.name} ===\n{content}")
                        except Exception:
                            pass
        
        elif ext == '.gz':
            with gzip.GzipFile(fileobj=BytesIO(body), mode='rb') as gf:
                content = gf.read().decode('utf-8', errors='replace')
                texts.append(content)
        
        return '\n'.join(texts)
    
    except Exception as e:
        log.warning("Failed to extract text from archive %s: %s", filename, e)
        return ""
    
    return ""


def _handle_text_upload(flow: http.HTTPFlow, filename: str, body: bytes, mime: str) -> None:
    """Handle text file upload by chunking and sending to QueueManager."""
    try:
        text = _extract_text_from_file(body, filename, mime)
    except Exception as e:
        log.warning("Failed to extract text from %s: %s", filename, e)
        return
    
    if not text.strip():
        log.debug("Empty text content in %s", filename)
        return

    # Generate message ID
    message_id = f"browser_{id(flow)}_{int(time.time())}"
    
    # Chunk the text
    chunks = chunk_text(text, _cfg.chunk_size_words, _cfg.chunk_overlap_words)
    if not chunks:
        return
    
    timestamp = datetime.datetime.utcnow().isoformat() + "Z"
    overall_decision = "ALLOW"
    
    log.info("Sending %d chunks from %s to QueueManager", len(chunks), filename)
    
    # Send each chunk to QueueManager and collect decisions
    for i, chunk_content in enumerate(chunks):
        payload = ChunkPayload(
            channel="browser",
            priority=False,
            message_id=message_id,
            chunk_id=i,
            total_chunks=len(chunks),
            content=chunk_content,
            word_count=len(chunk_content.split()),
            source_url=flow.request.pretty_url,
            filename=filename,
            timestamp=timestamp,
        )
        
        try:
            log.info("Sending chunk %d/%d to pipe %s", i + 1, len(chunks), _cfg.pipe_name)
            decision = pipe_client.send_and_receive(
                payload.to_dict(),
                _cfg.pipe_name,
                _cfg.timeout_seconds,
            )
            log.info("Received decision: %s for chunk %d/%d", decision, i + 1, len(chunks))

            # Streaming: if any chunk is BLOCK, overall is BLOCK
            if decision == "BLOCK":
                overall_decision = "BLOCK"
                log.info("BLOCK | chunk %d/%d | %s", i + 1, len(chunks), filename)
                break  # No need to send more chunks

            log.debug("ALLOW | chunk %d/%d | %s", i + 1, len(chunks), filename)

        except Exception as e:
            log.warning("Failed to send chunk %d: %s → fail_%s", i + 1, e, _cfg.fail_behavior)
            if not _cfg.fail_open():
                overall_decision = "BLOCK"
                break

    # Record decision
    with _browser_decisions_lock:
        _browser_decisions[message_id] = overall_decision
    
    # Apply decision
    if overall_decision == "BLOCK":
        _cache_blocked_url(flow.request.pretty_url)
        log.info("BLOCK | %s | %d bytes | %s", filename, len(body), flow.request.pretty_url)
        flow.response = http.Response.make(403, b"Upload blocked by DLP policy.", {"Content-Type": "text/plain"})
    else:
        log.info("ALLOW | %s | %d bytes | %s", filename, len(body), flow.request.pretty_url)


def _handle_binary_upload(flow: http.HTTPFlow, filename: str, body: bytes, mime: str) -> None:
    """Handle binary file upload using legacy temp file approach."""
    try:
        temp_path = _write_temp_file(body, filename)
    except OSError as e:
        log.error("Failed to write temp file for '%s': %s → fail_%s", filename, e, _cfg.fail_behavior)
        if not _cfg.fail_open():
            flow.kill()
        return

    payload = {
        "temp_path": temp_path,
        "url": flow.request.pretty_url,
        "method": flow.request.method,
        "content_type": flow.request.headers.get("content-type", ""),
        "effective_mime": mime,
        "filename": filename,
        "size_bytes": len(body),
        "timestamp": datetime.datetime.utcnow().isoformat() + "Z",
    }

    try:
        decision = pipe_client.send_and_receive(
            payload,
            _cfg.pipe_name,
            _cfg.timeout_seconds,
        )
        consumer_received = True
    except Exception as e:
        log.warning("Pipe error: %s → fail_%s", e, _cfg.fail_behavior)
        decision = "ALLOW" if _cfg.fail_open() else "BLOCK"
        consumer_received = False

    # Cleanup temp file if consumer didn't receive it
    if not consumer_received:
        _delete_temp_file(temp_path)

    if decision == "BLOCK":
        _cache_blocked_url(flow.request.pretty_url)
        log.info("BLOCK | %s | %d bytes | %s", filename, len(body), flow.request.pretty_url)
        flow.response = http.Response.make(403, b"Upload blocked by DLP policy.", {"Content-Type": "text/plain"})
    else:
        log.info("ALLOW | %s | %d bytes | %s", filename, len(body), flow.request.pretty_url)


# ---------------------------------------------------------------------------
# Helper functions (kept from original)
# ---------------------------------------------------------------------------

def _track_resumable_upload_initiation(flow: http.HTTPFlow) -> None:
    """Track Google Drive resumable upload initiation to get filename."""
    # The POST request to initiate resumable upload contains filename in headers
    title = flow.request.headers.get("x-upload-title", "")
    content_type = flow.request.headers.get("x-upload-content-type", "")
    
    if not title:
        # Try to get from query params
        title = flow.request.query.get("title", "")
    
    if not title:
        # Try to extract from URL
        title = _extract_filename_from_url(flow.request.path)
    
    if title:
        # Store for later PUT requests
        _pending_resumable[id(flow)] = title
        log.debug("Resumable upload initiation: tracking filename=%s", title)
    
    # Let the request pass through (it's just metadata)
    log.info("Resumable upload session initiated: %s", title or "unknown")


def _handle_resumable_upload(flow: http.HTTPFlow) -> None:
    """Handle Google Drive resumable upload (uploadType=resumable)."""
    url = flow.request.pretty_url
    
    # Check if already blocked
    if _is_blocked_url(url):
        log.debug("BLOCK (cached) | %s", url[:80])
        flow.response = http.Response.make(403, b"Upload blocked by DLP policy.", {"Content-Type": "text/plain"})
        return
    
    body = flow.request.content
    content_type = flow.request.headers.get("content-type", "").lower()
    
    # Skip metadata-only requests (no body or very small body)
    if not body or len(body) < _cfg.min_upload_size_bytes:
        log.debug("SKIP resumable: no body or too small (%d bytes)", len(body) if body else 0)
        return
    
    # Get filename from upload_id or URL
    upload_id = flow.request.query.get("upload_id", "")
    filename = ""
    
    if upload_id:
        filename = _resumable_filenames.get(upload_id, "")
        log.debug("Resumable upload: found filename=%s for upload_id", filename)
    
    if not filename:
        # Try to get filename from pending resumable tracking
        filename = _pending_resumable.get(id(flow), "")
        if filename:
            _resumable_filenames[upload_id] = filename
            log.debug("Resumable upload: tracked filename=%s", filename)
    
    if not filename:
        # Fallback: extract from URL
        filename = _extract_filename_from_url(flow.request.path)
        if not filename:
            filename = f"resumable_{id(flow)}_{int(time.time())}.bin"
        log.info("Resumable upload: using fallback filename=%s", filename)
    
    file_mime = content_type.split(";")[0].strip().lower() if content_type else ""
    
    log.info("Resumable upload detected: %s | filename=%s | size=%d bytes | mime=%s",
             url[:100], filename, len(body), file_mime)
    
    if not _matches_type_filter(filename, file_mime):
        log.debug("SKIP (type filter) | %s | %s", filename, file_mime)
        return
    
    # Handle text files: extract, chunk, send to QueueManager
    if _is_text_file(filename, file_mime):
        _handle_text_upload_flow(flow, filename, body, file_mime, url)
        return
    
    # Binary files: use temp file approach
    _handle_binary_upload_flow(flow, filename, body, file_mime, url)


def _handle_text_upload_flow(flow: http.HTTPFlow, filename: str, body: bytes, mime: str, url: str) -> None:
    """Handle text file upload by chunking and sending to QueueManager."""
    try:
        text = _extract_text_from_file(body, filename, mime)
    except Exception as e:
        log.warning("Failed to extract text from %s: %s", filename, e)
        return
    
    if not text.strip():
        log.debug("Empty text content in %s", filename)
        return
    
    # Generate message ID
    message_id = f"browser_{id(flow)}_{int(time.time())}"
    
    # Chunk the text
    chunks = chunk_text(text, _cfg.chunk_size_words, _cfg.chunk_overlap_words)
    if not chunks:
        return
    
    timestamp = datetime.datetime.utcnow().isoformat() + "Z"
    overall_decision = "ALLOW"
    
    log.info("Sending %d chunks from %s to QueueManager", len(chunks), filename)
    
    # Send each chunk to QueueManager and collect decisions
    for i, chunk_content in enumerate(chunks):
        payload = ChunkPayload(
            channel="browser",
            priority=False,
            message_id=message_id,
            chunk_id=i,
            total_chunks=len(chunks),
            content=chunk_content,
            word_count=len(chunk_content.split()),
            source_url=url,
            filename=filename,
            timestamp=timestamp,
        )
        
        try:
            log.info("Sending chunk %d/%d to pipe %s", i + 1, len(chunks), _cfg.pipe_name)
            decision = pipe_client.send_and_receive(
                payload.to_dict(),
                _cfg.pipe_name,
                _cfg.timeout_seconds,
            )
            log.info("Received decision: %s for chunk %d/%d", decision, i + 1, len(chunks))
            
            # Streaming: if any chunk is BLOCK, overall is BLOCK
            if decision == "BLOCK":
                overall_decision = "BLOCK"
                log.info("BLOCK | chunk %d/%d | %s", i + 1, len(chunks), filename)
                break  # No need to send more chunks
            
            log.debug("ALLOW | chunk %d/%d | %s", i + 1, len(chunks), filename)
            
        except Exception as e:
            log.warning("Failed to send chunk %d: %s → fail_%s", i + 1, e, _cfg.fail_behavior)
            if not _cfg.fail_open():
                overall_decision = "BLOCK"
                break
    
    # Record decision
    with _browser_decisions_lock:
        _browser_decisions[message_id] = overall_decision
    
    # Apply decision
    if overall_decision == "BLOCK":
        _cache_blocked_url(url)
        log.info("BLOCK | %s | %d bytes | %s", filename, len(body), url)
        flow.response = http.Response.make(403, b"Upload blocked by DLP policy.", {"Content-Type": "text/plain"})
    else:
        log.info("ALLOW | %s | %d bytes | %s", filename, len(body), url)


def _handle_binary_upload_flow(flow: http.HTTPFlow, filename: str, body: bytes, mime: str, url: str) -> None:
    """Handle binary file upload using temp file approach."""
    try:
        temp_path = _write_temp_file(body, filename)
    except OSError as e:
        log.error("Failed to write temp file for '%s': %s → fail_%s", filename, e, _cfg.fail_behavior)
        if not _cfg.fail_open():
            flow.kill()
        return
    
    payload = {
        "temp_path": temp_path,
        "url": url,
        "method": flow.request.method,
        "content_type": flow.request.headers.get("content-type", ""),
        "effective_mime": mime,
        "filename": filename,
        "size_bytes": len(body),
        "timestamp": datetime.datetime.utcnow().isoformat() + "Z",
    }
    
    try:
        decision = pipe_client.send_and_receive(
            payload,
            _cfg.pipe_name,
            _cfg.timeout_seconds,
        )
        consumer_received = True
    except Exception as e:
        log.warning("Pipe error: %s → fail_%s", e, _cfg.fail_behavior)
        decision = "ALLOW" if _cfg.fail_open() else "BLOCK"
        consumer_received = False
    
    # Cleanup temp file if consumer didn't receive it
    if not consumer_received:
        _delete_temp_file(temp_path)
    
    if decision == "BLOCK":
        _cache_blocked_url(url)
        log.info("BLOCK | %s | %d bytes | %s", filename, len(body), url)
        flow.response = http.Response.make(403, b"Upload blocked by DLP policy.", {"Content-Type": "text/plain"})
    else:
        log.info("ALLOW | %s | %d bytes | %s", filename, len(body), url)


def _is_upload(flow: http.HTTPFlow) -> bool:
    if flow.request.method not in ("POST", "PUT"):
        return False

    host = flow.request.pretty_host.lower()
    if any(host == d or host.endswith("." + d) for d in _cfg.domain_blocklist):
        log.debug("SKIP (domain blocklist) | %s", host)
        return False

    content_type = flow.request.headers.get("content-type", "").lower()

    if "multipart/form-data" in content_type or "multipart/related" in content_type:
        return True

    body_len = len(flow.request.content or b"")
    if body_len < _cfg.min_upload_size_bytes:
        return False

    if not _has_upload_url_keyword(flow):
        return False

    if not _has_filename_signal(flow):
        return False

    return True


def _has_upload_url_keyword(flow: http.HTTPFlow) -> bool:
    url_lower = flow.request.path.lower()
    return any(kw in url_lower for kw in _cfg.upload_url_keywords)


def _has_filename_signal(flow: http.HTTPFlow) -> bool:
    cd = flow.request.headers.get("content-disposition", "").lower()
    if "filename" in cd:
        return True

    query = flow.request.query
    if "filename" in query or "file_name" in query or "file" in query:
        return True

    upload_id = query.get("upload_id", "")
    if upload_id and upload_id in _resumable_filenames:
        return True

    path_segment = flow.request.path.split("?")[0].rsplit("/", 1)[-1]
    ext = os.path.splitext(path_segment)[1].lower()
    if ext:
        if not _cfg.has_type_filter():
            return True
        if ext in _cfg.extensions:
            return True

    return False


def _matches_type_filter(filename: str, mime: str) -> bool:
    if not _cfg.has_type_filter():
        return True
    ext = os.path.splitext(filename)[1].lower()
    if ext and ext in _cfg.extensions:
        return True
    if mime and mime in _cfg.mime_types:
        return True
    return False


def _extract_filename_from_url(path: str) -> str:
    """Extract filename from URL path or query parameters."""
    from urllib.parse import urlparse, parse_qs
    
    # Parse query parameters
    parsed = urlparse(path)
    query = parse_qs(parsed.query)
    
    # Check for filename in query params
    for key in ["filename", "file_name", "name", "title"]:
        if key in query:
            return query[key][0]
    
    # Try to get filename from path segment
    path_segment = parsed.path.rstrip("/").rsplit("/", 1)[-1]
    if path_segment and "." in path_segment:
        return path_segment
    
    return ""


def _extract_filename(flow: http.HTTPFlow) -> str:
    upload_id = flow.request.query.get("upload_id", "")
    if upload_id:
        filename = _resumable_filenames.get(upload_id, "")
        if filename:
            return filename

    content_type = flow.request.headers.get("content-type", "")
    if "multipart/form-data" in content_type.lower():
        name = _filename_from_multipart(flow.request.content, content_type)
        if name:
            return name

    path = flow.request.path.split("?")[0].rstrip("/")
    segment = path.rsplit("/", 1)[-1]
    return segment if segment else "upload"


def _filename_from_multipart(body: bytes, content_type: str) -> str:
    try:
        boundary = None
        for part in content_type.split(";"):
            part = part.strip()
            if part.startswith("boundary="):
                boundary = part[len("boundary="):].strip('"')
                break
        if not boundary:
            return ""
        delimiter = ("--" + boundary).encode()
        for chunk in body.split(delimiter):
            if b"filename=" not in chunk:
                continue
            header_end = chunk.find(b"\r\n\r\n")
            if header_end == -1:
                continue
            raw_headers = chunk[:header_end].lstrip(b"\r\n")
            parser = BytesHeaderParser()
            msg = parser.parsebytes(b"Content-Disposition: " + _extract_cd(raw_headers))
            params = msg.get_params(header="content-disposition")
            for key, val in params:
                if key == "filename":
                    return val
    except Exception:
        pass
    return ""


def _extract_cd(raw_headers: bytes) -> bytes:
    for line in raw_headers.split(b"\r\n"):
        if line.lower().startswith(b"content-disposition:"):
            return line[len(b"content-disposition:"):].strip()
    return b""


def _parse_multipart_related(body: bytes, content_type: str):
    try:
        boundary = None
        for segment in content_type.split(";"):
            segment = segment.strip()
            if segment.lower().startswith("boundary="):
                boundary = segment[len("boundary="):].strip('"').strip("'")
                break
        if not boundary:
            log.warning("multipart/related: no boundary found")
            return "", body, ""

        log.debug("multipart/related: boundary=%s, body size=%d", boundary, len(body))

        delimiter = ("--" + boundary).encode()
        parts = []
        for chunk in body.split(delimiter):
            chunk = chunk.lstrip(b"\r\n")
            stripped = chunk.rstrip(b"\r\n")
            if not stripped or stripped == b"--":
                continue
            sep = b"\r\n\r\n" if b"\r\n\r\n" in chunk else b"\n\n"
            header_end = chunk.find(sep)
            if header_end == -1:
                continue
            part_body = chunk[header_end + len(sep):].rstrip(b"\r\n")
            parts.append(part_body)

        log.debug("multipart/related: found %d parts", len(parts))

        if len(parts) < 2:
            log.warning("multipart/related: expected >=2 parts, got %d", len(parts))
            return "", body, ""

        # Part 0: JSON metadata
        try:
            metadata = json.loads(parts[0].decode("utf-8"))
            log.debug("multipart/related metadata: %s", metadata.keys() if isinstance(metadata, dict) else "not a dict")
            filename = metadata.get("name") or metadata.get("title") or ""
        except Exception as e:
            log.debug("multipart/related: failed to parse metadata: %s", e)
            filename = ""

        # Part 1: file content
        file_body_bytes = parts[1]
        
        log.info("multipart/related: extracted filename=%s, content size=%d", filename, len(file_body_bytes))
        return filename, file_body_bytes, ""

    except Exception as e:
        log.warning("Failed to parse multipart/related: %s", e)
        return "", body, ""


def _track_resumable_initiation_batch(flow: http.HTTPFlow) -> None:
    try:
        content_type = flow.request.headers.get("content-type", "")
        boundary = None
        for segment in content_type.split(";"):
            segment = segment.strip()
            if segment.lower().startswith("boundary="):
                boundary = segment[len("boundary="):].strip('"').strip("'")
                break
        if not boundary:
            return

        body = flow.request.content
        if not body:
            return

        delimiter = ("--" + boundary).encode()
        for chunk in body.split(delimiter):
            chunk = chunk.lstrip(b"\r\n")
            stripped = chunk.rstrip(b"\r\n")
            if not stripped or stripped == b"--":
                continue
            sep = b"\r\n\r\n" if b"\r\n\r\n" in chunk else b"\n\n"
            header_end = chunk.find(sep)
            if header_end == -1:
                continue
            part_body = chunk[header_end + len(sep):]

            inner_sep = b"\r\n\r\n" if b"\r\n\r\n" in part_body else b"\n\n"
            inner_split = part_body.find(inner_sep)
            if inner_split == -1:
                continue
            json_bytes = part_body[inner_split + len(inner_sep):].rstrip(b"\r\n")
            if not json_bytes:
                continue

            metadata = json.loads(json_bytes.decode("utf-8"))
            filename = metadata.get("title") or metadata.get("name") or ""
            if filename:
                _pending_resumable[id(flow)] = filename
                log.debug("Batch resumable init: queued filename=%r", filename)
                
                # Also store by upload_id if available
                upload_id = flow.request.query.get("upload_id", "")
                if upload_id:
                    _resumable_filenames[upload_id] = filename
                    log.debug("Batch resumable: mapped upload_id=%s -> filename=%s", upload_id[:20], filename)
            return
    except Exception as e:
        log.debug("Batch resumable init parse failed: %s", e)


def _is_blocked_url(url: str) -> bool:
    with _blocked_url_cache_lock:
        expiry = _blocked_url_cache.get(url)
        if expiry is None:
            return False
        now = time.monotonic()
        if now >= expiry:
            del _blocked_url_cache[url]
            return False
        return True


def _cache_blocked_url(url: str) -> None:
    now = time.monotonic()
    with _blocked_url_cache_lock:
        expired_keys = [k for k, exp in _blocked_url_cache.items() if now >= exp]
        for k in expired_keys:
            del _blocked_url_cache[k]
        _blocked_url_cache[url] = now + _BLOCK_CACHE_TTL


def _write_temp_file(body: bytes, filename: str) -> str:
    tmp_dir = _cfg.resolved_temp_dir()
    base, ext = os.path.splitext(filename) if filename else ("upload", "")
    dest = os.path.join(tmp_dir, filename if filename else "upload")
    counter = 0
    while os.path.exists(dest):
        counter += 1
        dest = os.path.join(tmp_dir, f"{base}_{counter}{ext}")
    with open(dest, "wb") as f:
        f.write(body)
    return dest


def _delete_temp_file(path: str) -> None:
    try:
        os.unlink(path)
    except OSError as e:
        log.warning("Could not delete temp file %s: %s", path, e)
