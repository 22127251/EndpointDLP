"""
Text extraction from various file formats.

Plain-text extraction (extract_text)
-------------------------------------
Supported:
  Plaintext   : .txt .md .json .yaml .yml .csv .log  (and any unrecognised extension)
  MS Office   : .docx  .xlsx  .pptx
  OpenDocument: .odt  .ods  .odp
  PDF         : .pdf   (via PyMuPDF)

Tabular extraction (extract_tabular)
--------------------------------------
Returns a TabularData object with per-column structure so that the analysis
engine can use column-header context matching instead of character proximity.

Supported:
  .csv / .tsv          — first row assumed to be column headers
  .xlsx (Excel)        — first row of each sheet assumed to be headers
  .ods (OpenDocument)  — first row of each table assumed to be headers
  .docx                — embedded tables (first row = headers) + body paragraphs
  .odt                 — embedded tables (first row = headers) + body text
"""

from __future__ import annotations

import csv as _csv
import io
import zipfile
from dataclasses import dataclass, field
from pathlib import Path

# ---------------------------------------------------------------------------
# Errors
# ---------------------------------------------------------------------------

class ExtractionTooLarge(Exception):
    """Raised when the running count of extracted characters exceeds the
    configured ``max_chars`` cap. Streaming extractors raise this *during* the
    parse (before the full body is materialized) so a pathological file cannot
    OOM the process; the orchestrator maps it to the channel's failure_mode
    (reason=text_cap). Carries the char count seen at the abort point."""

    def __init__(self, char_count: int):
        self.char_count = char_count
        super().__init__(f"extracted text exceeded cap ({char_count} chars)")


# ---------------------------------------------------------------------------
# Tabular data types
# ---------------------------------------------------------------------------

@dataclass
class ColumnBlock:
    header: str
    values: list[str]       # one entry per data row; empty cells → ""
    sheet: str | None       # None for single-sheet formats (CSV/TSV)


@dataclass
class TabularData:
    columns: list[ColumnBlock]
    # Free-text paragraphs (prose) that are NOT part of any table. These are
    # analyzed with bounded character proximity (engine.analyze), not the
    # column/row context used for `columns`, because prose has no row/column
    # structure — resolving its context by "same column/row" over-credits
    # context to far-apart values. One entry per source paragraph.
    body: list[str] = field(default_factory=list)


# ---------------------------------------------------------------------------
# Routing helpers
# ---------------------------------------------------------------------------

# CSV is tabular (header'd columns). PDF is intentionally NOT here: PyMuPDF
# find_tables is ~45 s on large PDFs, so PDF is routed to fast plain-text
# extraction (page.get_text); recall is unaffected (every value is still
# detected/counted), only header-based context becomes proximity-based.
_TABULAR_SUFFIXES = {".csv", ".tsv", ".xlsx", ".ods", ".docx", ".odt"}


def is_tabular(file_path: str | Path) -> bool:
    """Return True if the file should be processed with extract_tabular()."""
    return Path(file_path).suffix.lower() in _TABULAR_SUFFIXES


# ---------------------------------------------------------------------------
# Plain-text extraction
# ---------------------------------------------------------------------------

def extract_text(file_path: str | Path, max_chars: int | None = None) -> str:
    """Return all readable text from *file_path* as a single string.

    If *max_chars* is set and the extracted text exceeds it, raise
    ExtractionTooLarge. For these formats the check is post-hoc: the input is
    already bounded by the orchestrator's max_file_bytes and none of them is the
    14×-expanding case. The office formats that *do* blow up (docx/odt/ods/xlsx)
    go through extract_tabular, which enforces the cap during the parse."""
    path = Path(file_path)
    suffix = path.suffix.lower()

    if suffix == ".docx":
        text = _extract_docx(path)
    elif suffix == ".xlsx":
        text = _extract_xlsx(path)
    elif suffix == ".pptx":
        text = _extract_pptx(path)
    elif suffix in {".odt", ".ods", ".odp"}:
        text = _extract_odf(path)
    elif suffix == ".pdf":
        text = _extract_pdf(path)
    else:
        text = _extract_plaintext(path)
    if max_chars is not None and len(text) > max_chars:
        raise ExtractionTooLarge(len(text))
    return text


def _extract_plaintext(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


def _extract_docx(path: Path) -> str:
    import docx

    doc = docx.Document(str(path))
    parts: list[str] = []
    for para in doc.paragraphs:
        if para.text:
            parts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text:
                    parts.append(cell.text)
    return "\n".join(parts)


def _extract_xlsx(path: Path) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    parts: list[str] = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            for cell in row:
                if cell is not None:
                    parts.append(str(cell))
    return "\n".join(parts)


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text:
                        parts.append(para.text)
    return "\n".join(parts)


def _extract_odf(path: Path) -> str:
    from odf import text as odf_text
    from odf.opendocument import load as odf_load
    from odf.element import Element

    doc = odf_load(str(path))
    parts: list[str] = []

    def _walk(node: Element) -> None:
        if node.nodeType == node.TEXT_NODE:
            v = node.data
            if v:
                parts.append(v)
        for child in node.childNodes:
            _walk(child)

    _walk(doc.body)
    return "\n".join(parts)


def _extract_pdf(path: Path) -> str:
    import fitz

    parts: list[str] = []
    with fitz.open(str(path)) as doc:
        for page in doc:
            page_text = page.get_text()
            if page_text:
                parts.append(page_text)
    return "\n".join(parts)


# ---------------------------------------------------------------------------
# Tabular extraction
# ---------------------------------------------------------------------------

def extract_tabular(file_path: str | Path, max_chars: int | None = None) -> TabularData:
    """Extract column-structured data from tabular and document files.

    If *max_chars* is set, extraction is refused once the running count of
    extracted characters exceeds it (ExtractionTooLarge). The streaming office
    readers (docx/odt/ods/xlsx) enforce it *during* the parse so a zip-expanding
    file cannot OOM the process before any post-hoc check; csv/pdf — bounded by
    max_file_bytes and not subject to zip expansion — are checked once after."""
    path = Path(file_path)
    suffix = path.suffix.lower()

    if suffix == ".csv":
        td = _extract_csv_tabular(path, delimiter=",")
    elif suffix == ".tsv":
        td = _extract_csv_tabular(path, delimiter="\t")
    elif suffix in (".xlsx", ".ods"):
        return _extract_calamine_tabular(path, max_chars)
    elif suffix == ".docx":
        return _extract_docx_tabular(path, max_chars)
    elif suffix == ".odt":
        return _extract_odt_tabular(path, max_chars)
    elif suffix == ".pdf":
        td = _extract_pdf_tabular(path)
    else:
        # Fallback: no table structure — treat the whole file as free-text body.
        text = extract_text(path, max_chars)
        return TabularData(columns=[], body=text.splitlines())
    _enforce_tabular_cap(td, max_chars)
    return td


def _enforce_tabular_cap(td: TabularData, max_chars: int | None) -> None:
    """Post-hoc extracted-text cap for the non-streaming tabular readers
    (csv/pdf): raise ExtractionTooLarge if the total extracted characters exceed
    *max_chars*. Counts cell values + body (headers are negligible), matching the
    running count the streaming readers accumulate."""
    if max_chars is None:
        return
    total = (sum(len(v) for c in td.columns for v in c.values)
             + sum(len(b) for b in td.body))
    if total > max_chars:
        raise ExtractionTooLarge(total)


# ---- CSV / TSV ----

def _extract_csv_tabular(path: Path, delimiter: str = ",") -> TabularData:
    with open(path, encoding="utf-8", errors="replace", newline="") as f:
        reader = _csv.reader(f, delimiter=delimiter)
        rows = list(reader)

    if not rows:
        return TabularData(columns=[])

    headers = [h.strip() for h in rows[0]]
    data_rows = rows[1:]

    # Fallback: if every non-empty header cell is numeric, treat file as headerless
    non_empty = [h for h in headers if h]
    if non_empty and all(_is_numeric_string(h) for h in non_empty):
        headers = [f"Column {i + 1}" for i in range(len(headers))]
        data_rows = rows

    columns: list[ColumnBlock] = []
    for col_idx, header in enumerate(headers):
        values = [
            row[col_idx].strip() if col_idx < len(row) else ""
            for row in data_rows
        ]
        columns.append(ColumnBlock(header=header, values=values, sheet=None))

    return TabularData(columns=columns)


# ---- Excel (.xlsx) + ODS (.ods) — python-calamine (Rust) ----

def _coerce_cell(v) -> str:
    """Coerce a python-calamine cell value to text.

    Text cells come back as ``str`` (so leading zeros and ``+84`` prefixes are
    preserved verbatim); numbers/dates are stringified. An *integral* float drops
    its ``.0`` so a number-typed ID cell becomes ``"12345"`` not ``"12345.0"``.
    ``bool`` is handled before the float/int branches (it is an ``int`` subclass).
    Empty cells are calamine's own ``""`` and pass straight through."""
    if isinstance(v, str):
        return v
    if isinstance(v, bool):
        return str(v)
    if isinstance(v, float) and v.is_integer():
        return str(int(v))
    return str(v)


def _extract_calamine_tabular(path: Path, max_chars: int | None = None) -> TabularData:
    """Extract column-structured data from .xlsx / .ods with python-calamine (Rust).

    Replaces the openpyxl (xlsx) and lxml-iterparse (ods) readers — calamine
    parses both formats several times faster (corpus: ~3 s → <1 s) and at lower
    memory. Rows are pulled one at a time via ``CalamineSheet.iter_rows`` so the
    running character count can raise ExtractionTooLarge *mid-parse*, before the
    whole grid is materialized in Python, keeping the Phase 5 cap that
    ``to_python()`` (a whole-grid materializer) could not. Each cell is coerced
    to text by ``_coerce_cell``; the row-major grid is handed to the shared
    ``_grid_to_columns`` (same header / numeric-header-fallback logic as before).

    Hard dependency, **no fallback**: if python-calamine is not importable the
    import below raises and extraction fails loudly, so it is never ambiguous
    which reader ran (decision #3 in the perf/failmode plan)."""
    from python_calamine import CalamineWorkbook

    columns: list[ColumnBlock] = []
    char_count = 0
    wb = CalamineWorkbook.from_path(str(path))
    try:
        for idx in range(len(wb.sheet_names)):
            sheet = wb.get_sheet_by_index(idx)
            grid: list[list[str]] = []
            for row in sheet.iter_rows():
                cells = [_coerce_cell(c) for c in row]
                grid.append(cells)
                char_count += sum(len(c) for c in cells)
                if max_chars is not None and char_count > max_chars:
                    raise ExtractionTooLarge(char_count)
            _grid_to_columns(grid, sheet.name, columns)
    finally:
        wb.close()

    return TabularData(columns=columns)


# ---- Shared: build ColumnBlocks from a row-major grid ----

def _grid_to_columns(grid: list[list[str]], sheet: str | None,
                     columns: list[ColumnBlock]) -> None:
    """Append columns for a row-major *grid* (row 0 = headers) to *columns*.
    Mirrors the legacy numeric-header fallback (all-numeric header → headerless)."""
    if not grid:
        return
    ncols = max(len(r) for r in grid)
    headers = grid[0]
    data_rows = grid[1:]

    non_empty = [h for h in headers if h.strip()]
    if non_empty and all(_is_numeric_string(h) for h in non_empty):
        headers = [f"Column {i + 1}" for i in range(ncols)]
        data_rows = grid

    for ci in range(ncols):
        header = headers[ci].strip() if ci < len(headers) else ""
        values = [(r[ci] if ci < len(r) else "") for r in data_rows]
        columns.append(ColumnBlock(header=header, values=values, sheet=sheet))


# ---- DOCX (Word document — tables + body paragraphs) — fast lxml ----

def _extract_docx_tabular(path: Path, max_chars: int | None = None) -> TabularData:
    """Stream word/document.xml with lxml.iterparse (C parser, element-by-element
    memory release) instead of building the whole tree with etree.fromstring —
    bounds extraction memory for huge documents and provides the early-abort hook
    (mirrors the ODS streamer). Paragraphs inside table cells are excluded from
    the body column (else their PII is double-counted). Handles simple
    (non-nested) tables, which covers typical documents.

    A running character count is accumulated across body paragraphs and table
    cell values; if *max_chars* is set and the count exceeds it, raise
    ExtractionTooLarge mid-parse so a pathological file is refused before its
    full text is materialized."""
    from lxml import etree

    W = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"

    with zipfile.ZipFile(path) as z:
        data = z.read("word/document.xml")

    def _cell_text(tc) -> str:
        return " ".join("".join(p.itertext()) for p in tc.iter(W + "p")).strip()

    columns: list[ColumnBlock] = []
    body: list[str] = []
    char_count = 0
    tbl_idx = 0
    tbl_depth = 0  # >0 while inside a (possibly nested) table; body p's are skipped

    for event, el in etree.iterparse(io.BytesIO(data), events=("start", "end")):
        tag = el.tag
        if event == "start":
            if tag == W + "tbl":
                tbl_depth += 1
            continue

        if tag == W + "tbl":
            tbl_depth -= 1
            if tbl_depth != 0:  # inner table; outer end processes the whole tree
                continue
            sheet = f"Table {tbl_idx + 1}"
            tbl_idx += 1
            rows = list(el.iter(W + "tr"))
            if rows:
                raw_headers = [_cell_text(c) for c in rows[0].iter(W + "tc")]
                for ci, header in enumerate(raw_headers):
                    values: list[str] = []
                    for tr in rows[1:]:
                        cells = list(tr.iter(W + "tc"))
                        v = _cell_text(cells[ci]) if ci < len(cells) else ""
                        values.append(v)
                        char_count += len(v)
                    columns.append(ColumnBlock(header=header.strip(), values=values, sheet=sheet))
                if max_chars is not None and char_count > max_chars:
                    raise ExtractionTooLarge(char_count)
            el.clear()
        elif tag == W + "p" and tbl_depth == 0:
            # Body paragraph (w:p not inside a table). Word does not nest
            # paragraphs, so clearing on end is safe.
            t = "".join(el.itertext()).strip()
            if t:
                body.append(t)
                char_count += len(t)
                if max_chars is not None and char_count > max_chars:
                    raise ExtractionTooLarge(char_count)
            el.clear()

    return TabularData(columns=columns, body=body)


# ---- ODT (OpenDocument Text — tables + body text) — fast lxml ----

def _extract_odt_tabular(path: Path, max_chars: int | None = None) -> TabularData:
    """Stream content.xml with lxml.iterparse instead of building the whole tree
    with etree.fromstring — bounds extraction memory and provides the early-abort
    hook (mirrors the ODS streamer). Paragraphs inside table cells are excluded
    from the body column. Simple (non-nested) tables.

    Accumulates a running character count over body paragraphs and table cell
    values; if *max_chars* is set and the count exceeds it, raise
    ExtractionTooLarge mid-parse."""
    from lxml import etree

    T = "{urn:oasis:names:tc:opendocument:xmlns:table:1.0}"
    TX = "{urn:oasis:names:tc:opendocument:xmlns:text:1.0}"
    CELL_TAGS = (T + "table-cell", T + "covered-table-cell")

    with zipfile.ZipFile(path) as z:
        data = z.read("content.xml")

    def _cell_text(cell) -> str:
        return " ".join("".join(p.itertext()) for p in cell.iter(TX + "p")).strip()

    columns: list[ColumnBlock] = []
    body: list[str] = []
    char_count = 0
    tbl_idx = 0
    tbl_depth = 0  # >0 while inside a (possibly nested) table; body p's are skipped

    for event, el in etree.iterparse(io.BytesIO(data), events=("start", "end")):
        tag = el.tag
        if event == "start":
            if tag == T + "table":
                tbl_depth += 1
            continue

        if tag == T + "table":
            tbl_depth -= 1
            if tbl_depth != 0:  # inner table; outer end processes the whole tree
                continue
            sheet = f"Table {tbl_idx + 1}"
            tbl_idx += 1
            rows = list(el.iter(T + "table-row"))
            if rows:
                header_cells = [c for c in rows[0] if c.tag in CELL_TAGS]
                raw_headers = [_cell_text(c) for c in header_cells]
                for ci, header in enumerate(raw_headers):
                    values: list[str] = []
                    for tr in rows[1:]:
                        cells = [c for c in tr if c.tag in CELL_TAGS]
                        v = _cell_text(cells[ci]) if ci < len(cells) else ""
                        values.append(v)
                        char_count += len(v)
                    columns.append(ColumnBlock(header=header.strip(), values=values, sheet=sheet))
                if max_chars is not None and char_count > max_chars:
                    raise ExtractionTooLarge(char_count)
            el.clear()
        elif tag == TX + "p" and tbl_depth == 0:
            # Body paragraph (text:p not inside a table). ODF paragraphs do not
            # nest, so clearing on end is safe.
            t = "".join(el.itertext()).strip()
            if t:
                body.append(t)
                char_count += len(t)
                if max_chars is not None and char_count > max_chars:
                    raise ExtractionTooLarge(char_count)
            el.clear()

    return TabularData(columns=columns, body=body)


# ---- PDF (PyMuPDF — tables + body text) ----

def _extract_pdf_tabular(path: Path) -> TabularData:
    import fitz

    columns: list[ColumnBlock] = []
    body_texts: list[str] = []
    # Indices into `columns` and column count of the most recently completed
    # multi-row table.  Kept outside the page loop so they survive page turns
    # and can be used to merge 1-row stubs that PyMuPDF creates when a table
    # row wraps across a page boundary.
    last_col_indices: list[int] = []
    last_col_count: int = 0

    with fitz.open(str(path)) as doc:
        for page_num, page in enumerate(doc):
            finder = page.find_tables()
            table_bboxes = [t.bbox for t in finder.tables]

            for table_idx, table in enumerate(finder.tables):
                rows = table.extract()
                if not rows:
                    continue

                if len(rows) == 1:
                    # PyMuPDF splits a page-spanning table into a normal table on
                    # page N and a 1-row "stub" table on page N+1.  The normal
                    # code path would treat that single row as a header with no
                    # data rows, silently discarding every PII value in it.
                    # Instead, merge the stub back into the previous table.
                    stub_row = rows[0]
                    if last_col_indices and len(stub_row) == last_col_count:
                        non_empty = sum(
                            1 for c in stub_row if c is not None and str(c).strip()
                        )
                        if non_empty >= (last_col_count + 1) // 2:
                            # Most cells are populated → the full last row of the
                            # previous table overflowed to this page.  Append each
                            # cell as a new data value in its matching column.
                            for ci, block_idx in enumerate(last_col_indices):
                                cell = stub_row[ci] if ci < len(stub_row) else None
                                columns[block_idx].values.append(
                                    str(cell) if cell is not None else ""
                                )
                        else:
                            # Only a minority of cells have content → just the tail
                            # of one or more cell values overflowed.  Concatenate
                            # each non-empty stub cell onto the last value of the
                            # matching column so the regex sees the complete number
                            # (e.g. "4111 5555 6666" + "\n7777" → full 16-digit VISA).
                            for ci, block_idx in enumerate(last_col_indices):
                                cell = stub_row[ci] if ci < len(stub_row) else None
                                cell_str = str(cell).strip() if cell is not None else ""
                                if cell_str:
                                    if columns[block_idx].values:
                                        columns[block_idx].values[-1] += "\n" + cell_str
                                    else:
                                        columns[block_idx].values.append(cell_str)
                    continue  # stub handled; skip normal header/data processing

                sheet = f"Page {page_num + 1} Table {table_idx + 1}"
                raw_headers = [str(c) if c is not None else "" for c in rows[0]]
                start_idx = len(columns)
                for col_idx, header in enumerate(raw_headers):
                    values: list[str] = []
                    for row in rows[1:]:
                        cell = row[col_idx] if col_idx < len(row) else None
                        values.append(str(cell) if cell is not None else "")
                    columns.append(ColumnBlock(header=header, values=values, sheet=sheet))
                # Record this table's column range so the next stub can find it.
                last_col_indices = list(range(start_idx, len(columns)))
                last_col_count = len(raw_headers)

            # Collect body text at word granularity instead of block granularity.
            # page.get_text("blocks") can merge a paragraph with an adjacent table
            # into one oversized block; filtering at block level then discards the
            # paragraph.  Individual words have tight bboxes so the overlap check
            # correctly distinguishes table words from paragraph words even when
            # PyMuPDF assigned them the same block_no.
            #
            # word_data: (block_no, line_no) → [(x0, word_text)]
            #   block_no – PyMuPDF block index on this page
            #   line_no  – top-to-bottom line counter within that block
            #   x0       – left edge of the word (used for left-to-right ordering)
            word_data: dict[tuple[int, int], list[tuple[float, str]]] = {}
            for word_info in page.get_text("words"):
                wx0, wy0, wx1, wy1, word_text, block_no, line_no, _ = word_info
                in_table = any(
                    not (wx1 <= tb[0] or tb[2] <= wx0 or wy1 <= tb[1] or tb[3] <= wy0)
                    for tb in table_bboxes
                )
                if not in_table:
                    word_data.setdefault((block_no, line_no), []).append((wx0, word_text))

            block_lines: dict[int, list[tuple[int, str]]] = {}
            for (block_no, line_no), words in word_data.items():
                words.sort(key=lambda w: w[0])  # left-to-right within a line
                block_lines.setdefault(block_no, []).append(
                    (line_no, " ".join(w[1] for w in words))
                )

            # Each block becomes one body_texts entry so that context words and
            # PII values in the same paragraph stay in the same "cell" for the
            # engine's inline-context matching.
            for block_no in sorted(block_lines.keys()):
                lines = sorted(block_lines[block_no], key=lambda x: x[0])
                block_text = "\n".join(line for _, line in lines).strip()
                if block_text:
                    body_texts.append(block_text)

    return TabularData(columns=columns, body=body_texts)


# ---------------------------------------------------------------------------
# Utility
# ---------------------------------------------------------------------------

def _is_numeric_string(s: str) -> bool:
    return s.strip().lstrip("-").replace(".", "", 1).isdigit()
